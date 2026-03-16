from pathlib import Path
from pptx import Presentation
import json

BASE_DIR = Path(__file__).resolve().parent.parent
PPTX_PATH = BASE_DIR / "input" / "sample_report.pptx"
OUTPUT_DIR = BASE_DIR / "output" / "data"
OUTPUT_JSON = OUTPUT_DIR / "slides.json"


def extract_notes_text(slide) -> str:
    if not slide.has_notes_slide:
        return ""

    text = slide.notes_slide.notes_text_frame.text or ""
    return text.strip()


def main():
    if not PPTX_PATH.exists():
        print(f"PowerPoint file not found: {PPTX_PATH}")
        return

    presentation = Presentation(PPTX_PATH)
    slides_data = []

    for i, slide in enumerate(presentation.slides, start=1):
        notes = extract_notes_text(slide)

        slide_data = {
            "slide_number": i,
            "image_path": f"/assets/slides/slide_{i}.png",
            "avatar_script": notes
        }
        slides_data.append(slide_data)

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    with open(OUTPUT_JSON, "w", encoding="utf-8") as f:
        json.dump(slides_data, f, indent=2, ensure_ascii=False)

    print(f"slides.json generated successfully at: {OUTPUT_JSON}")


if __name__ == "__main__":
    main()