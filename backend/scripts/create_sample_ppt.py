from pptx import Presentation
import os

# Create presentation
prs = Presentation()

slides_content = [
    ("Market Overview", "Welcome to the Q3 market landscape report."),
    ("Growth Trends", "This slide explains how the market has grown over the past year."),
    ("Future Outlook", "Here we discuss expected trends for the next quarter.")
]

for title, notes in slides_content:
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    slide.shapes.title.text = title
    slide.placeholders[1].text = "Sample content for the slide"

    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes

# Save file
output_path = "../input/sample_report.pptx"
os.makedirs("../input", exist_ok=True)
prs.save(output_path)

print("Sample PowerPoint created successfully!")